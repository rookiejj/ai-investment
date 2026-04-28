#!/usr/bin/env python3
"""
브리픽 카드사 심사용 PPTX 자동 생성 스크립트.

스크린샷은 docs/screenshots/photo_2026-04-28 15.41.*.jpeg 파일을 순서대로 사용.
사용: python3 scripts/build_audit_pptx.py
산출물: docs/브리픽 카드사 심사용 자료.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image  # python-pptx가 PIL 의존
import sys

ROOT = Path(__file__).resolve().parents[1]
SS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "브리픽 카드사 심사용 자료.pptx"

# 슬라이드 사이즈: 16:9
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# 색상
BRIEFICK_GREEN = RGBColor(0x18, 0xE2, 0x99)
TEXT_DARK = RGBColor(0x0D, 0x0D, 0x0D)
TEXT_MUT = RGBColor(0x66, 0x66, 0x66)

# 스크린샷 매핑
SCREENS = {
    "main":      SS / "photo_2026-04-28 15.41.10.jpeg",
    "footer":    SS / "photo_2026-04-28 15.41.12.jpeg",
    "modal":     SS / "photo_2026-04-28 15.41.14.jpeg",
    "confirm":   SS / "photo_2026-04-28 15.41.16.jpeg",
    "pg":        SS / "photo_2026-04-28 15.41.17.jpeg",
    "card_auth": SS / "photo_2026-04-28 15.41.20.jpeg",
    "complete":  SS / "photo_2026-04-28 15.41.21.jpeg",
}


def add_title(slide, text, size=Pt(28), bold=True, color=TEXT_DARK,
              left=Inches(0.6), top=Inches(0.4), width=Inches(12.1), height=Inches(0.7)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return box


def add_subtitle(slide, text, top=Inches(1.05), color=TEXT_MUT, size=Pt(14)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.5))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.name = "Pretendard"
    return box


def add_image_centered(slide, path: Path, top=Inches(1.65), max_h=Inches(5.5), max_w=Inches(11.5)):
    if not path.exists():
        print(f"[warn] missing: {path}", file=sys.stderr)
        return None
    img = Image.open(path)
    iw, ih = img.size
    ratio = iw / ih
    target_w = max_w
    target_h = Emu(int(target_w / ratio))
    if target_h > max_h:
        target_h = max_h
        target_w = Emu(int(target_h * ratio))
    left = Emu(int((SLIDE_W - target_w) / 2))
    pic = slide.shapes.add_picture(str(path), left, top, width=target_w, height=target_h)
    # 살짝 그림자 느낌의 보더 박스 (단순 사각형)
    return pic


def add_footer_pagenum(slide, idx, total):
    box = slide.shapes.add_textbox(Inches(11.8), Inches(7.05), Inches(1.3), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{idx} / {total}"
    run.font.size = Pt(10)
    run.font.color.rgb = TEXT_MUT
    run.font.name = "Pretendard"


def add_brand_badge(slide, top=Inches(7.05)):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(3.0), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "브리픽 · briefick"
    run.font.size = Pt(10)
    run.font.color.rgb = BRIEFICK_GREEN
    run.font.bold = True
    run.font.name = "Pretendard"


def cover_slide(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 배경 강조
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    bg.line.fill.background()
    # 큰 타이틀
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.0))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "브리픽 (Briefick)\n"
    r1.font.size = Pt(56)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_DARK
    r1.font.name = "Pretendard"
    r2 = p.add_run()
    r2.text = "카드사 심사용 자료"
    r2.font.size = Pt(36)
    r2.font.color.rgb = BRIEFICK_GREEN
    r2.font.bold = True
    r2.font.name = "Pretendard"
    # 서브타이틀
    sub = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = "매일 경제 뉴스 카카오톡 친구톡 구독 서비스 — 결제 화면 흐름"
    srun.font.size = Pt(18)
    srun.font.color.rgb = TEXT_MUT
    srun.font.name = "Pretendard"
    # URL/회사 정보
    info = s.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.6))
    itf = info.text_frame
    ip = itf.paragraphs[0]
    irun = ip.add_run()
    irun.text = "https://roysbriefing.vercel.app    ·    로이소프트    ·    사업자 737-42-01455"
    irun.font.size = Pt(12)
    irun.font.color.rgb = TEXT_MUT
    irun.font.name = "Pretendard"
    return s


def screen_slide(prs, idx_text, title, subtitle, screenshot_key, page_idx, total_pages):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 상단 인덱스 + 타이틀
    add_title(s, f"{idx_text}. {title}")
    if subtitle:
        add_subtitle(s, subtitle)
    add_image_centered(s, SCREENS[screenshot_key])
    add_brand_badge(s)
    add_footer_pagenum(s, page_idx, total_pages)
    return s


def section_slide(prs, label, page_idx, total_pages):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    # 회색 배경
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x0D)
    bg.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = BRIEFICK_GREEN
    r.font.name = "Pretendard"
    return s


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 슬라이드 정의 (인덱스, 제목, 부제, 스크린샷 키)
    plan = [
        ("1", "메인 화면",
         "https://roysbriefing.vercel.app — 일반 인증 결제 진입점 (헤더 우측 '구독하기' 버튼)",
         "main"),
        ("2", "홈페이지 하단 정보",
         "사업자번호 · 상호 · 대표자 · 주소 · 일반전화번호 · 이메일 노출",
         "footer"),
        ("3", "결제 상품 선택 — 구독 모달",
         "헤더 '구독하기' 클릭 → 전화번호 입력 + 광고성 정보 수신 동의 체크",
         "modal"),
        ("4", "결제하기 — 번호 확인",
         "오타 방지를 위한 번호 큰 글씨 재확인 단계 · 기존 구독자에겐 연장 후 만료일 안내",
         "confirm"),
        ("5", "PG사 결제창",
         "포트원 V2 + 갤럭시아머니트리 결제창 노출 (카드사 선택)",
         "pg"),
        ("6", "카드사 인증",
         "선택된 카드사 인증/결제 화면",
         "card_auth"),
        ("7", "결제 완료",
         "결제 완료 후 구독 연장 안내 + 카카오 채널 친구 추가 강조",
         "complete"),
    ]

    total_pages = 1 + len(plan)  # cover + steps

    cover_slide(prs)

    for i, (idx_text, title, subtitle, key) in enumerate(plan, start=1):
        screen_slide(prs, idx_text, title, subtitle, key, page_idx=i + 1, total_pages=total_pages + 1)

    # 보충 슬라이드: 자동결제 / 빌링 미사용 안내
    s = prs.slide_layouts[6]
    note = prs.slides.add_slide(s)
    add_title(note, "참고 — 정기(빌링·구독) 결제 흐름")
    add_subtitle(note, "본 서비스는 자동 결제(빌링) 미운영 — 만료 시 동일 결제 흐름으로 수동 재결제 (재구독 페이지: /renew)")
    # 본문 박스
    body = note.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.0))
    btf = body.text_frame
    btf.word_wrap = True
    items = [
        "• 매월 자동 결제 미운영. 매 월 만료일 전 사용자가 직접 재결제하는 구조.",
        "• 만료 D-1 카카오 친구톡으로 재구독 안내 발송 (https://roysbriefing.vercel.app/renew 링크 포함).",
        "• 재구독 페이지의 결제 흐름은 위 3~7단계와 100% 동일 (구독 모달 → 번호 확인 → PG사 결제창 → 카드사 인증 → 결제 완료).",
        "• 결제 1건당 30일 이용권 부여, 다음 결제 없이는 만료 후 자동 비활성화.",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(16)
        run.font.color.rgb = TEXT_DARK
        run.font.name = "Pretendard"
        p.space_after = Pt(8)
    add_brand_badge(note)
    add_footer_pagenum(note, total_pages + 1, total_pages + 1)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    main()
